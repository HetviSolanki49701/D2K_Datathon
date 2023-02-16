from pptx import Presentation
from pptx.util import Inches

# create a new presentation
prs = Presentation()

# add a title slide
title_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Summary Of the Research Paper"
subtitle.text = "Automated Generation made by Team MangoDB"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 32, 96)


# add a bullet slide
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Bullet Slide"
tf = body_shape.text_frame
tf.text = "Bullet 1"
p = tf.add_paragraph()
p.text = "Bullet 2"
p.level = 1

# add an image slide
img_path = "image.jpg"
pic_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(pic_slide_layout)
left = top = Inches(1)
width = height = Inches(5.5)
pic = slide.shapes.add_picture(img_path, left, top, width, height)

# save the presentation
prs.save("test1.pptx")
