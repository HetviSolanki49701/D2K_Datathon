import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

# create a new presentation
prs = Presentation()

chartTitle = ["Number of Sentences Before Summary", "Number of Sentences After Summary"]
chartTitle2 = ["Number of Words Before Summary", "Number of Words After Summary"]
words = [300, 100]
sentences = [3000, 300]


# add a title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Research Paper Summary"
subtitle.text = "Curated by Team MangoDB"
# title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
# subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 32, 96)

title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)



chart_data = CategoryChartData()
chart_data.categories = chartTitle
chart_data.add_series('Analysis of Sentences', (words))
x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
slide1.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

chart_data = CategoryChartData()
chart_data.categories = chartTitle2
chart_data.add_series('Analysis of Words', (sentences))
x, y, cx, cy = Inches(.5), Inches(.5), Inches(9), Inches(6.5)
slide2.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)


# # add a bullet slide
# bullet_slide_layout = prs.slide_layouts[1]
# slide = prs.slides.add_slide(bullet_slide_layout)
# shapes = slide.shapes
# title_shape = shapes.title
# body_shape = shapes.placeholders[1]
# title_shape.text = "Bullet Slide"
# tf = body_shape.text_frame
# tf.text = "Bullet 1"
# p = tf.add_paragraph()
# p.text = "Bullet 2"
# p.level = 1

# # add an image slide
# img_path = 'image.jpg'
# pic_slide_layout = prs.slide_layouts[6]
# slide = prs.slides.add_slide(pic_slide_layout)
# left = top = Inches(1)
# width = height = Inches(5.5)
# pic = slide.shapes.add_picture(img_path, left, top, width, height)

# save the presentation
prs.save('test3.pptx')
